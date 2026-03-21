"""
File processor for various document formats
"""
import PyPDF2
from pathlib import Path
from typing import Optional


class FileProcessor:
    """Process various file formats"""
    
    async def process(self, file_path: Path) -> Optional[str]:
        """Process file and extract text"""
        ext = file_path.suffix.lower()
        
        processors = {
            ".pdf": self._process_pdf,
            ".txt": self._process_text,
            ".docx": self._process_docx,
            ".doc": self._process_doc,
            ".pptx": self._process_pptx,
            ".xlsx": self._process_xlsx,
            ".csv": self._process_csv,
            ".json": self._process_json,
            ".md": self._process_markdown,
        }
        
        processor = processors.get(ext)
        if not processor:
            return None
        
        return processor(file_path)
    
    def _process_pdf(self, file_path: Path) -> str:
        """Extract text from PDF"""
        text = ""
        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text()
        except Exception as e:
            print(f"Error processing PDF: {e}")
        return text
    
    def _process_text(self, file_path: Path) -> str:
        """Read text file"""
        try:
            return file_path.read_text(encoding="utf-8")
        except Exception:
            try:
                return file_path.read_text(encoding="latin-1")
            except Exception as e:
                print(f"Error reading text file: {e}")
                return ""
    
    def _process_docx(self, file_path: Path) -> str:
        """Extract text from DOCX"""
        try:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            return text
        except Exception as e:
            print(f"Error processing DOCX: {e}")
            return ""
    
    def _process_doc(self, file_path: Path) -> str:
        """Extract text from DOC (requires python-docx)"""
        # For old .doc format, would need python-docx with workaround
        return self._process_docx(file_path)
    
    def _process_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Error processing PPTX: {e}")
            return ""
    
    def _process_xlsx(self, file_path: Path) -> str:
        """Extract text from XLSX"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text += f"Sheet: {sheet}\n"
                for row in ws.iter_rows(values_only=True):
                    text += " | ".join(str(cell) if cell else "" for cell in row) + "\n"
            return text
        except Exception as e:
            print(f"Error processing XLSX: {e}")
            return ""
    
    def _process_csv(self, file_path: Path) -> str:
        """Read CSV file"""
        try:
            return file_path.read_text(encoding="utf-8")
        except Exception as e:
            print(f"Error reading CSV: {e}")
            return ""
    
    def _process_json(self, file_path: Path) -> str:
        """Read JSON file"""
        try:
            import json
            with open(file_path, "r") as f:
                data = json.load(f)
            return json.dumps(data, indent=2)
        except Exception as e:
            print(f"Error reading JSON: {e}")
            return ""
    
    def _process_markdown(self, file_path: Path) -> str:
        """Read Markdown file"""
        try:
            return file_path.read_text(encoding="utf-8")
        except Exception as e:
            print(f"Error reading Markdown: {e}")
            return ""
